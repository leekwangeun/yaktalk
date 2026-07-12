# 발표자료(아키텍처 섹션) 생성 — python scripts/build_slides.py
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "presentation" / "약톡_아키텍처.pptx"
OUT.parent.mkdir(exist_ok=True)

FONT = "Malgun Gothic"
DARK = "05343E"; TEAL = "028090"; PURPLE = "6C5CE7"; GRAY = "566B75"
MINT = "02C39A"; INK = "16323B"; MUTE = "5F7178"; LIGHT = "F5FAFB"

KIND = {  # fill, border, title color, sub color
    "data":    ("E2F0F0", TEAL,   "05343E", "2C6E72"),
    "model":   ("EAE7FB", PURPLE, "2C2568", "5B51B0"),
    "infra":   ("EBEEF0", GRAY,   "263238", "556670"),
    "neutral": ("FFFFFF", "C9D6DB", INK,    MUTE),
    "result":  ("FEF3E0", "E8912D","7A4A08", "9A6212"),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 13.333, 7.5
blank = prs.slide_layouts[6]


def rgb(h): return RGBColor.from_string(h)


def add_bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = rgb(color); r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2, r._element)
    return r


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs[0], tuple): runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for t, sz, col, bold in para:
            r = p.add_run(); r.text = t
            r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = rgb(col)
    return tb


def box(slide, x, y, w, h, title, sub=None, kind="neutral", tsize=14, ssize=10.5):
    fill, border, tc, sc = KIND[kind]
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    sp.line.color.rgb = rgb(border); sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = 0.11
    except Exception: pass
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(5); tf.margin_top = tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(tsize)
    r.font.bold = True; r.font.color.rgb = rgb(tc)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub; r2.font.name = FONT; r2.font.size = Pt(ssize)
        r2.font.color.rgb = rgb(sc)
    return sp


def arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = rgb(color); c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return c


def circle_num(slide, x, y, d, n, color):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = rgb(color); o.line.fill.background()
    o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.name = FONT; r.font.size = Pt(13)
    r.font.bold = True; r.font.color.rgb = rgb("FFFFFF")
    return o


def header(slide, num, title, sub):
    circle_num(slide, 0.6, 0.52, 0.5, num, TEAL)
    text(slide, 1.2, 0.44, 9.5, 0.5, [(title, 25, INK, True)])
    text(slide, 1.2, 0.98, 11.4, 0.4, [(sub, 13, MUTE, False)])


def legend(slide, x, y, items):
    cx = x
    for label, color in items:
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(y + 0.03), Inches(0.16), Inches(0.16))
        d.fill.solid(); d.fill.fore_color.rgb = rgb(color); d.line.fill.background(); d.shadow.inherit = False
        lw = 0.14 * len(label) + 0.12
        text(slide, cx + 0.24, y - 0.04, lw, 0.3, [(label, 11, MUTE, False)], wrap=False)
        cx += 0.24 + lw + 0.4


# ---------------- Slide 1: 섹션 표지 ----------------
s = prs.slides.add_slide(blank)
add_bg(s, DARK)
text(s, 0.9, 1.15, 11, 0.5, [("약물 상호작용 확인 챗봇  ·  약톡", 15, MINT, True)])
text(s, 0.85, 2.35, 12, 1.4, [("시스템 아키텍처", 56, "FFFFFF", True)])
text(s, 0.9, 3.75, 11.5, 0.6, [("공공데이터 수집부터 직접 학습한 AI, 그리고 배포까지", 18, "CFE6E6", False)])
steps = [("구축 과정", "식약처 데이터 → 모델 학습"), ("시스템 구조", "Vercel + HF Spaces"),
         ("질문 처리 흐름", "입력 → 판정 → 응답")]
bx = 0.9
for i, (t, d) in enumerate(steps, 1):
    circle_num(s, bx, 5.25, 0.46, i, MINT)
    text(s, bx + 0.6, 5.22, 3.3, 0.4, [(t, 15, "FFFFFF", True)], wrap=False)
    text(s, bx + 0.6, 5.62, 3.3, 0.4, [(d, 11, "9FC6C6", False)], wrap=False)
    bx += 4.05

# ---------------- Slide 2: 구축 과정 ----------------
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT)
header(s, 1, "구축 과정", "식약처 공공데이터가 판단 근거이자 AI 학습 데이터가 되는 이중 활용 구조")
stages = [
    ("식약처 공공데이터", "허가정보·DUR·e약은요\n4개 API", "data"),
    ("수집 · 전처리", "20만 행 → 정제\ndrug_light.db 43MB", "data"),
    ("학습 데이터 합성", "질문 4만 문장\n응답 9천 쌍", "purple_special"),
    ("Colab 파인튜닝", "KoELECTRA ×2\nKoGPT2 ×1", "model"),
    ("로컬 모델 자산", "외부 API 0\n오프라인 동작", "model"),
]
n = len(stages); bw = 2.15; gap = (W - 1.4 - bw * n) / (n - 1)
y = 2.5; bh = 1.55
xs = []
for i, (t, sub, kind) in enumerate(stages):
    x = 0.7 + i * (bw + gap); xs.append(x)
    k = "data" if kind == "data" else "model"
    b = box(s, x, y, bw, bh, t, sub.replace("\n", "  "), k, tsize=13.5, ssize=10)
    if i < n - 1:
        arrow(s, x + bw, y + bh / 2, x + bw + gap, y + bh / 2, GRAY)
# 강조 캡션
cap = box(s, 0.7, 4.7, W - 1.4, 1.05,
          "핵심: 허가정보 제품명 사전 × 질문 템플릿 → 라벨링 학습 데이터를 직접 합성",
          "라벨 데이터가 없다는 한계를 '공공데이터로 데이터를 만드는' 방식으로 해결 — 같은 데이터가 판정 근거이자 학습 재료",
          "neutral", tsize=15, ssize=11.5)
cap.line.color.rgb = rgb(MINT); cap.line.width = Pt(1.5)
legend(s, 0.7, 6.15, [("식약처 데이터", TEAL), ("직접 학습 모델", PURPLE)])

# ---------------- Slide 3: 시스템 구조 ----------------
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT)
header(s, 2, "시스템 구조", "프론트(Vercel)와 백엔드(Hugging Face Spaces)를 분리 — 모든 모델·데이터는 컨테이너 내부에 탑재")
box(s, 0.7, 1.85, 3.0, 0.95, "사용자 브라우저", "채팅 · 노인/임산부 토글", "neutral", 14, 10.5)
arrow(s, 3.7, 2.32, 4.35, 2.32, GRAY)
box(s, 4.35, 1.85, 3.0, 0.95, "Vercel", "web/ 정적 채팅 UI", "neutral", 14, 10.5)
arrow(s, 5.85, 2.8, 5.85, 3.25, GRAY)
# HF 컨테이너
cont = slide = None
c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.3), Inches(W - 1.4), Inches(3.05))
c.fill.solid(); c.fill.fore_color.rgb = rgb("ECF1F2"); c.line.color.rgb = rgb(GRAY); c.line.width = Pt(1.25)
c.shadow.inherit = False
try: c.adjustments[0] = 0.04
except Exception: pass
text(s, 1.0, 3.42, 11, 0.4, [("Hugging Face Spaces  —  Docker · FastAPI(api.py) · 포트 7860", 13.5, "37474F", True)])
mods = [("① 자연어 이해", "사전 + 자모 + KoELECTRA", "model", "KoELECTRA ×2", "NER · 의도  55MB×2", "model"),
        ("② 위험도 판정", "DUR 검사 6종 · 엔티티", "data", "drug_light.db", "식약처 12테이블 43MB", "data"),
        ("③ 응답 생성", "KoGPT2 → 검증 게이트", "model", "KoGPT2 + 템플릿", "생성 481MB · 폴백", "model")]
mw = 3.6; mgap = (W - 1.4 - 0.6 - mw * 3) / 2; my = 3.95
for i, (t, sub, k, at, asub, ak) in enumerate(mods):
    x = 1.0 + i * (mw + mgap)
    box(s, x, my, mw, 1.0, t, sub, k, 14, 10.5)
    arrow(s, x + mw / 2, my + 1.0, x + mw / 2, my + 1.32, GRAY, 1.4)
    box(s, x, my + 1.35, mw, 0.92, at, asub, ak, 12.5, 10)
text(s, 0.7, 6.55, 8, 0.4,
     [[("오프라인 데모: ", 12.5, INK, True), ("uvicorn api:app --app-dir src  한 줄로 동일 구성 로컬 실행 (발표장 무중단)", 12.5, MUTE, False)]])
legend(s, 9.0, 6.55, [("직접 학습 모델", PURPLE), ("식약처 데이터", TEAL)])

# ---------------- Slide 4: 질문 처리 흐름 ----------------
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT)
header(s, 3, "질문 처리 흐름", '입력 "타이레놀이랑 판콜 같이 먹어도 돼?" 가 응답이 되기까지')
# 입력 콜아웃
ib = box(s, 0.7, 1.8, 3.1, 1.0, "사용자 입력", '"타이레놀이랑 판콜\n같이 먹어도 돼?"', "neutral", 13.5, 11)
steps = [("① 의도 분류", "KoELECTRA → 병용", "model"),
         ("② 약 이름 추출", "사전 + NER 앙상블", "model"),
         ("③ 제품→성분 변환", "허가정보 DB 조회", "data"),
         ("④ 위험도 검사", "DUR 병용금기·중복", "data"),
         ("⑤ 응답 생성", "KoGPT2 + 게이트", "model")]
n = len(steps); bw = 2.28; gap = (W - 1.4 - bw * n) / (n - 1); y = 3.35; bh = 1.5
for i, (t, sub, k) in enumerate(steps):
    x = 0.7 + i * (bw + gap)
    box(s, x, y, bw, bh, t, sub, k, 13, 10)
    if i < n - 1:
        arrow(s, x + bw, y + bh / 2, x + bw + gap, y + bh / 2, GRAY, 1.5)
arrow(s, 2.25, 2.8, 2.25, 3.35, GRAY, 1.5)
# 되묻기 분기 + 결과
br = box(s, 0.7, 5.25, 5.7, 1.05, "모호한 입력이면 → 되묻기",
         '"감기약이신가요? 판콜에스·판피린큐 중에…" (멀티턴, 문맥 유지)', "neutral", 13, 10.5)
res = box(s, 6.9, 5.25, W - 1.4 - 6.2, 1.05, "🟠 주의  —  아세트아미노펜 중복",
          "두 약 모두 아세트아미노펜 함유 → 과량 위험 (1일 최대 4,000mg)", "result", 14, 10.5)
legend(s, 0.7, 6.65, [("직접 학습 모델", PURPLE), ("식약처 데이터", TEAL)])

# ---------------- Slide 5: 설계 원칙 ----------------
s = prs.slides.add_slide(blank)
add_bg(s, DARK)
text(s, 0.85, 0.9, 12, 0.9, [("핵심 설계 원칙", 34, "FFFFFF", True)])
text(s, 0.9, 1.85, 11.5, 0.6, [('"AI가 창작할 수 있는 자리를 구조적으로 없앴다"', 18, MINT, True)])
cols = [("이해", "직접 학습한 모델", "사용자의 자연어·상품명·오타를\nKoELECTRA NER과 사전으로 해석", PURPLE),
        ("판단", "100% 식약처 공공데이터", "위험 여부는 DUR 규칙만으로 결정\nAI는 의학적 판단에 관여하지 않음", TEAL),
        ("표현", "검증 게이트 + 템플릿", "KoGPT2 생성문이 사실 검증을\n통과할 때만 사용, 실패 시 템플릿", PURPLE)]
cw = 3.75; cgap = (W - 1.7 - cw * 3) / 2; cy = 2.95; ch = 2.75
for i, (tag, title, desc, col) in enumerate(cols):
    x = 0.85 + i * (cw + cgap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(cy), Inches(cw), Inches(ch))
    card.fill.solid(); card.fill.fore_color.rgb = rgb("0C4650")
    card.line.color.rgb = rgb(col); card.line.width = Pt(1.5); card.shadow.inherit = False
    try: card.adjustments[0] = 0.06
    except Exception: pass
    circle_num(s, x + 0.35, cy + 0.35, 0.5, i + 1, col)
    text(s, x + 0.95, cy + 0.34, cw - 1.1, 0.5, [(tag, 19, "FFFFFF", True)])
    text(s, x + 0.35, cy + 1.15, cw - 0.7, 0.5, [(title, 14.5, MINT, True)])
    text(s, x + 0.35, cy + 1.7, cw - 0.7, 1.0,
         [(desc.split("\n")[0], 12, "CFE6E6", False)] if "\n" not in desc else
         [[(desc.split("\n")[0], 12, "CFE6E6", False)], [(desc.split("\n")[1], 12, "CFE6E6", False)]])
text(s, 0.85, 6.15, 12, 0.6,
     [[("결과: ", 13, MINT, True),
       ("외부 AI API 의존 0 · 완전 오프라인 동작 · 근거는 항상 식약처 원문 인용", 13, "CFE6E6", False)]])

prs.save(str(OUT))
print("saved:", OUT)
